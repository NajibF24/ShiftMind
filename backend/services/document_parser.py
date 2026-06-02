"""
Document Parser — Multi-format document parsing and text chunking.
Supports: PDF, DOCX, XLSX, PPTX, TXT, MD
"""
import os
import io
import logging
from typing import Optional

logger = logging.getLogger(__name__)


def parse_document(file_content: bytes, filename: str) -> Optional[str]:
    """
    Parse document bytes into plain text based on file extension.
    Returns extracted text or None on failure.
    """
    ext = os.path.splitext(filename)[1].lower()

    try:
        if ext == ".pdf":
            return _parse_pdf(file_content)
        elif ext in (".docx", ".doc"):
            return _parse_docx(file_content)
        elif ext == ".xlsx":
            return _parse_xlsx(file_content)
        elif ext == ".pptx":
            return _parse_pptx(file_content)
        elif ext in (".txt", ".md"):
            return _parse_text(file_content)
        else:
            logger.warning(f"Unsupported file extension: {ext} for file {filename}")
            return None
    except Exception as e:
        logger.error(f"Error parsing {filename}: {e}")
        return None


def _parse_pdf(content: bytes) -> str:
    """Extract text from PDF using PyPDF2."""
    from PyPDF2 import PdfReader

    reader = PdfReader(io.BytesIO(content))
    texts = []
    for page_num, page in enumerate(reader.pages, 1):
        text = page.extract_text()
        if text and text.strip():
            texts.append(f"[Halaman {page_num}]\n{text.strip()}")
    return "\n\n".join(texts)


def _parse_docx(content: bytes) -> str:
    """Extract text from DOCX using python-docx."""
    from docx import Document

    doc = Document(io.BytesIO(content))
    texts = []

    for para in doc.paragraphs:
        if para.text.strip():
            # Preserve heading structure
            if para.style and para.style.name and para.style.name.startswith("Heading"):
                level = para.style.name.replace("Heading ", "").replace("Heading", "1")
                try:
                    prefix = "#" * int(level)
                except ValueError:
                    prefix = "#"
                texts.append(f"{prefix} {para.text.strip()}")
            else:
                texts.append(para.text.strip())

    # Also extract tables
    for table in doc.tables:
        table_text = _extract_table(table)
        if table_text:
            texts.append(table_text)

    return "\n\n".join(texts)


def _extract_table(table) -> str:
    """Extract text from a docx table."""
    rows = []
    for row in table.rows:
        cells = [cell.text.strip() for cell in row.cells]
        rows.append(" | ".join(cells))
    if rows:
        return "Tabel:\n" + "\n".join(rows)
    return ""


def _parse_xlsx(content: bytes) -> str:
    """Extract text from XLSX using openpyxl."""
    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(content), read_only=True, data_only=True)
    texts = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        sheet_text = f"Sheet: {sheet_name}\n"
        rows_data = []

        for row in ws.iter_rows(values_only=True):
            cell_values = [str(c) if c is not None else "" for c in row]
            if any(v.strip() for v in cell_values):
                rows_data.append(" | ".join(cell_values))

        if rows_data:
            sheet_text += "\n".join(rows_data)
            texts.append(sheet_text)

    wb.close()
    return "\n\n".join(texts)


def _parse_pptx(content: bytes) -> str:
    """Extract text from PPTX using python-pptx."""
    from pptx import Presentation

    prs = Presentation(io.BytesIO(content))
    texts = []

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_texts = [f"[Slide {slide_num}]"]
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_texts.append(shape.text.strip())
        if len(slide_texts) > 1:  # More than just the header
            texts.append("\n".join(slide_texts))

    return "\n\n".join(texts)


def _parse_text(content: bytes) -> str:
    """Parse plain text / markdown files."""
    # Try UTF-8 first, then fall back to latin-1
    try:
        return content.decode("utf-8")
    except UnicodeDecodeError:
        return content.decode("latin-1")


def chunk_text(text: str, chunk_size: int = 500, overlap: int = 50) -> list[str]:
    """
    Split text into chunks of approximately `chunk_size` words with `overlap` word overlap.
    Tries to split at paragraph/sentence boundaries for cleaner chunks.
    """
    if not text or not text.strip():
        return []

    # Split into paragraphs first
    paragraphs = [p.strip() for p in text.split("\n\n") if p.strip()]

    chunks = []
    current_chunk_words = []
    current_word_count = 0

    for para in paragraphs:
        para_words = para.split()
        para_word_count = len(para_words)

        # If a single paragraph exceeds chunk_size, split it by sentences
        if para_word_count > chunk_size:
            # Split by sentences (periods, or newlines)
            sentences = []
            for line in para.split("\n"):
                for s in line.replace(". ", ".\n").split("\n"):
                    if s.strip():
                        sentences.append(s.strip())

            for sentence in sentences:
                sent_words = sentence.split()
                if current_word_count + len(sent_words) > chunk_size and current_chunk_words:
                    # Save current chunk
                    chunks.append(" ".join(current_chunk_words))
                    # Keep overlap words
                    current_chunk_words = current_chunk_words[-overlap:] if overlap > 0 else []
                    current_word_count = len(current_chunk_words)

                current_chunk_words.extend(sent_words)
                current_word_count += len(sent_words)
        else:
            # Check if adding this paragraph would exceed limit
            if current_word_count + para_word_count > chunk_size and current_chunk_words:
                chunks.append(" ".join(current_chunk_words))
                current_chunk_words = current_chunk_words[-overlap:] if overlap > 0 else []
                current_word_count = len(current_chunk_words)

            current_chunk_words.extend(para_words)
            current_word_count += para_word_count

    # Don't forget the last chunk
    if current_chunk_words:
        chunks.append(" ".join(current_chunk_words))

    return chunks
